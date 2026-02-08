#!/usr/bin/env python3
"""
Generate a visually polished PowerPoint presentation.
Standalone script - not integrated with the project codebase.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
FIGURES_DIR = ROOT / "figures"
OUT_PATH = ROOT / "presentation.pptx"

# ── Design Tokens ──────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)
DEEP_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
BLUE      = RGBColor(0x1E, 0x88, 0xE5)
LIGHT_BLUE= RGBColor(0xBB, 0xDE, 0xFB)
ICE_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
TEAL      = RGBColor(0x00, 0x96, 0x88)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
CHARCOAL  = RGBColor(0x21, 0x21, 0x21)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)
MID_GRAY  = RGBColor(0x75, 0x75, 0x75)
LIGHT_GRAY= RGBColor(0xBD, 0xBD, 0xBD)
ACCENT_ORANGE = RGBColor(0xFF, 0x6F, 0x00)
ACCENT_GREEN  = RGBColor(0x2E, 0x7D, 0x32)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def add_line(slide, left, top, width, color, thickness=Pt(2)):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.height = int(thickness)
    return shape


def tb(slide, left, top, width, height, text, size=18, bold=False,
       color=CHARCOAL, align=PP_ALIGN.LEFT, font="Calibri", spacing_after=None,
       line_spacing=None):
    """Add a single-paragraph text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    if spacing_after is not None:
        p.space_after = spacing_after
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return tf


def bullets(slide, left, top, width, height, items, size=16, color=CHARCOAL,
            space=Pt(6), font="Calibri"):
    """
    Add bulleted text. items = list of (text, bold, level).
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (txt, bld, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = bld
        p.font.color.rgb = color
        p.font.name = font
        p.level = lvl
        p.space_after = space
    return tf


def img(slide, path, left, top, max_w, max_h):
    """Add image fitted within bounds."""
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    aspect = iw / ih
    if max_w / aspect <= max_h:
        w, h = max_w, max_w / aspect
    else:
        h, w = max_h, max_h * aspect
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))
    return w, h


def img_centered(slide, path, top, max_w=11.0, max_h=5.0):
    """Add image centered horizontally."""
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    aspect = iw / ih
    if max_w / aspect <= max_h:
        w, h = max_w, max_w / aspect
    else:
        h, w = max_h, max_h * aspect
    left = (13.333 - w) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))
    return w, h


def slide_header(slide, title, subtitle=None):
    """Standard slide header with accent bar."""
    # Top accent bar
    add_shape(slide, 0, 0, SLIDE_W, Inches(0.06), BLUE)
    # Title
    tb(slide, 0.8, 0.35, 11.5, 0.7, title, size=30, bold=True, color=NAVY, font="Calibri Light")
    # Underline
    add_line(slide, Inches(0.8), Inches(0.95), Inches(1.8), BLUE, thickness=Pt(3))
    if subtitle:
        tb(slide, 0.8, 1.05, 11.5, 0.4, subtitle, size=14, color=MID_GRAY)


def card(slide, left, top, width, height, bg_color=ICE_BLUE, border_color=None):
    """Add a rounded-ish card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ── Slides ─────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    # Decorative accent bar on left
    add_shape(s, 0, 0, Inches(0.15), SLIDE_H, BLUE)

    # Subtle horizontal line
    add_shape(s, Inches(2), Inches(3.4), Inches(9), Pt(1), RGBColor(0x2A, 0x3F, 0x6A))

    tb(s, 1.5, 1.2, 10.5, 1.5,
       "Optimal Oyster Reef\nSite Selection",
       size=44, bold=True, color=WHITE, font="Calibri Light", line_spacing=Pt(52))

    tb(s, 1.5, 3.6, 10.0, 0.8,
       "Combining ODE Modeling, Heuristic Optimization,\nand Mixed-Integer Quadratic Programming",
       size=20, color=LIGHT_BLUE, font="Calibri Light", line_spacing=Pt(28))

    # Subtle bottom bar
    add_shape(s, 0, Inches(6.6), SLIDE_W, Pt(1), RGBColor(0x2A, 0x3F, 0x6A))

    tb(s, 1.5, 5.8, 5, 0.5, "Paul Marouf", size=22, bold=True, color=WHITE)
    tb(s, 1.5, 6.3, 5, 0.4, "Chesapeake Bay Oyster Restoration Project", size=14, color=LIGHT_BLUE)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "The Problem")

    # Left column - stat cards
    card(s, 0.8, 1.5, 3.6, 1.6, bg_color=RGBColor(0xFF, 0xEB, 0xEE))
    tb(s, 1.1, 1.65, 3.0, 0.6, "99%", size=48, bold=True, color=RGBColor(0xC6, 0x28, 0x28), align=PP_ALIGN.CENTER)
    tb(s, 1.1, 2.35, 3.0, 0.5, "of Chesapeake Bay oysters lost\n10 billion  ->  < 100 million", size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    card(s, 4.8, 1.5, 3.6, 1.6, bg_color=ICE_BLUE)
    tb(s, 5.1, 1.65, 3.0, 0.6, "25 / 48", size=48, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)
    tb(s, 5.1, 2.35, 3.0, 0.5, "sites we can restore\nout of 48 candidates", size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    card(s, 8.8, 1.5, 3.6, 1.6, bg_color=RGBColor(0xE8, 0xF5, 0xE9))
    tb(s, 9.1, 1.65, 3.0, 0.6, "1.4T", size=48, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    tb(s, 9.1, 2.35, 3.0, 0.5, "possible combinations\nC(48, 25) - brute force impossible", size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Bottom - the question
    card(s, 0.8, 3.6, 11.7, 1.8, bg_color=WHITE, border_color=BLUE)
    tb(s, 1.2, 3.75, 11.0, 0.5,
       "The Optimization Question", size=20, bold=True, color=NAVY)
    tb(s, 1.2, 4.3, 11.0, 0.9,
       "Which 25 reef sites should be restored to maximize total adult oyster biomass,\n"
       "given the larval connectivity network between sites?",
       size=17, color=DARK_GRAY, line_spacing=Pt(24))

    # Key challenge
    tb(s, 0.8, 5.7, 11.7, 1.2,
       "Challenge: Each site's value depends on which other sites are also restored.\n"
       "Larvae flow between reefs through ocean currents, creating network effects\n"
       "that make this a nonlinear combinatorial problem.",
       size=14, color=MID_GRAY, line_spacing=Pt(20))


def slide_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "The Biological Model", "JARS ODE System  |  Provided by Dr. Leah Shaw")

    # Four state cards
    states = [
        ("J", "Juveniles", "Newly settled\nlarvae", RGBColor(0xE3, 0xF2, 0xFD)),
        ("A", "Adults", "Biomass\n(our objective)", RGBColor(0xE8, 0xF5, 0xE9)),
        ("R", "Reef / Shell", "Habitat\nstructure", RGBColor(0xFF, 0xF3, 0xE0)),
        ("S", "Sediment", "Smothering\nfactor", RGBColor(0xFC, 0xE4, 0xEC)),
    ]
    for i, (letter, name, desc, bg) in enumerate(states):
        x = 0.8 + i * 3.1
        card(s, x, 1.5, 2.7, 1.8, bg_color=bg)
        tb(s, x + 0.2, 1.6, 2.3, 0.7, letter, size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tb(s, x + 0.2, 2.2, 2.3, 0.3, name, size=14, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        tb(s, x + 0.2, 2.55, 2.3, 0.5, desc, size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Key equations / concepts
    card(s, 0.8, 3.7, 5.6, 3.0, bg_color=WHITE, border_color=LIGHT_GRAY)
    tb(s, 1.1, 3.85, 5.0, 0.4, "How Sites Interact", size=18, bold=True, color=NAVY)
    bullets(s, 1.1, 4.35, 5.0, 2.2, [
        ("Larval input = (P0 + P1' |A|^1.72) * L * f", False, 0),
        ("P0 = external larvae (from outside the network)", False, 0),
        ("P1 = internal connectivity matrix (56 x 56)", False, 0),
        ("|A|^1.72 = nonlinear adult biomass scaling", False, 0),
        ("L, f = habitat quality and survival functions", False, 0),
    ], size=13, color=DARK_GRAY, space=Pt(4))

    card(s, 6.8, 3.7, 5.7, 3.0, bg_color=WHITE, border_color=LIGHT_GRAY)
    tb(s, 7.1, 3.85, 5.0, 0.4, "Computational Cost", size=18, bold=True, color=NAVY)
    bullets(s, 7.1, 4.35, 5.0, 2.2, [
        ("Each evaluation = solve 4N-dim ODE to t=1000", False, 0),
        ("N = number of sites in the subset (up to 48)", False, 0),
        ("One ODE solve takes ~0.1-0.2 seconds", False, 0),
        ("Greedy needs ~900 evaluations", False, 0),
        ("Need smart methods, not just brute force", True, 0),
    ], size=13, color=DARK_GRAY, space=Pt(4))


def slide_connectivity(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Larval Connectivity Network")

    img(s, FIGURES_DIR / "fig8_connectivity_heatmap.png", 0.5, 1.3, 7.0, 5.8)

    card(s, 7.8, 1.5, 4.8, 5.0, bg_color=WHITE, border_color=LIGHT_GRAY)
    tb(s, 8.1, 1.65, 4.2, 0.4, "What This Shows", size=18, bold=True, color=NAVY)
    bullets(s, 8.1, 2.15, 4.2, 4.0, [
        ("56-site connectivity matrix from\noceanographic larval transport modeling", False, 0),
        ("Hot cells = strong larval flow\nbetween source and destination reefs", False, 0),
        ("Visible clustering: sites 10, 27-29,\n31-32, 40-41 form a connected hub", False, 0),
        ("Off-diagonal structure reveals\nlong-range larval pathways", False, 0),
        ("This network structure drives which\nsites are most valuable to restore", True, 0),
    ], size=13, color=DARK_GRAY, space=Pt(10))


def slide_contribution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "My Contribution: Three Solution Approaches")

    # Card 1: Heuristic
    card(s, 0.8, 1.4, 3.7, 5.2, bg_color=WHITE, border_color=BLUE)
    add_shape(s, Inches(0.8), Inches(1.4), Inches(3.7), Inches(0.55), BLUE)
    tb(s, 0.8, 1.42, 3.7, 0.5, "Heuristic Search", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(s, 1.1, 2.1, 3.2, 4.0, [
        ("Uses full ODE model", True, 0),
        ("Forward Greedy", True, 0),
        ("Add best site one at a time\nacross 25 rounds", False, 1),
        ("Backward Elimination", True, 0),
        ("Start with all 48 sites,\nremove least impactful", False, 1),
        ("Local Search (1-swap)", True, 0),
        ("Refine by swapping sites\nin and out of the set", False, 1),
    ], size=12, color=DARK_GRAY, space=Pt(4))

    # Card 2: MIQP
    card(s, 4.8, 1.4, 3.7, 5.2, bg_color=WHITE, border_color=ACCENT_ORANGE)
    add_shape(s, Inches(4.8), Inches(1.4), Inches(3.7), Inches(0.55), ACCENT_ORANGE)
    tb(s, 4.8, 1.42, 3.7, 0.5, "MIQP Surrogate", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(s, 5.1, 2.1, 3.2, 4.0, [
        ("Key innovation", True, 0),
        ("Quadratic Objective", True, 0),
        ("max  sum Pe_i x_i\n   + sum W_ij x_i x_j", False, 1),
        ("Surrogate Weights", True, 0),
        ("W_ij = P1_ij * (A*)^1.72\napprox. network effect", False, 1),
        ("Exact Solution", True, 0),
        ("Gurobi solver: provably\noptimal in 0.32 seconds", False, 1),
    ], size=12, color=DARK_GRAY, space=Pt(4))

    # Card 3: Constrained
    card(s, 8.8, 1.4, 3.7, 5.2, bg_color=WHITE, border_color=ACCENT_GREEN)
    add_shape(s, Inches(8.8), Inches(1.4), Inches(3.7), Inches(0.55), ACCENT_GREEN)
    tb(s, 8.8, 1.42, 3.7, 0.5, "Constrained MIQP", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bullets(s, 9.1, 2.1, 3.2, 4.0, [
        ("Real-world extensions", True, 0),
        ("Community Constraints", True, 0),
        ("Minimum reefs per\ngeographic region for fairness", False, 1),
        ("Reef Sizing", True, 0),
        ("Optimize area allocation\nper site under total budget", False, 1),
        ("Combined Model", True, 0),
        ("Communities + sizing\nin a single formulation", False, 1),
    ], size=12, color=DARK_GRAY, space=Pt(4))


def slide_greedy_curve(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Greedy Selection: Marginal Returns Curve")

    img_centered(s, FIGURES_DIR / "fig7_greedy_curve.png", 1.2, max_w=10.5, max_h=5.2)

    tb(s, 0.8, 6.6, 11.7, 0.5,
       "Each site adds biomass, with a jump at sites 17-18 (hub sites 40, 41). "
       "Greedy+Local ultimately achieves the best ODE-validated score of 1.862.",
       size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_scores(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "ODE-Validated Performance Comparison")

    img_centered(s, FIGURES_DIR / "fig1_score_comparison.png", 1.2, max_w=9.5, max_h=5.0)

    tb(s, 0.8, 6.5, 11.7, 0.6,
       "All solutions validated through the full nonlinear JARS ODE (tmax=1000, realistic P0). "
       "Greedy+Local achieves highest biomass. MIQP is 8% behind but solves 6,000x faster.",
       size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_timing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Computational Time: Speed vs Quality")

    img(s, FIGURES_DIR / "fig9_timing.png", 0.5, 1.3, 7.0, 5.0)

    card(s, 7.8, 1.5, 4.8, 2.2, bg_color=WHITE, border_color=LIGHT_GRAY)
    tb(s, 8.1, 1.65, 4.2, 0.4, "The Trade-off", size=18, bold=True, color=NAVY)
    bullets(s, 8.1, 2.15, 4.2, 1.3, [
        ("Greedy+Local: 32 min, best score", False, 0),
        ("Backward: 4.7 min, 2.3% gap", False, 0),
        ("MIQP: 0.32 sec, 8% gap", False, 0),
        ("MIQP is 6,000x faster", True, 0),
    ], size=13, color=DARK_GRAY, space=Pt(6))

    card(s, 7.8, 4.0, 4.8, 2.3, bg_color=ICE_BLUE, border_color=BLUE)
    tb(s, 8.1, 4.15, 4.2, 0.4, "Why MIQP Matters", size=18, bold=True, color=NAVY)
    bullets(s, 8.1, 4.65, 4.2, 1.3, [
        ("Enables rapid what-if analysis", False, 0),
        ("Test 100s of constraint scenarios", False, 0),
        ("Interactive decision support", False, 0),
        ("Provably optimal (under surrogate)", True, 0),
    ], size=13, color=DARK_GRAY, space=Pt(6))


def slide_gap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Optimality Gap Analysis")

    img_centered(s, FIGURES_DIR / "fig10_optimality_gap.png", 1.2, max_w=10.0, max_h=5.0)

    tb(s, 0.8, 6.5, 11.7, 0.6,
       "The 8% MIQP gap reflects the surrogate approximation error, not the solver. "
       "Community constraints cost ~26% biomass - quantifying the price of geographic fairness.",
       size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_overlap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Method Agreement: Which Sites Overlap?")

    img(s, FIGURES_DIR / "fig2_site_overlap_heatmap.png", 0.5, 1.3, 6.5, 5.8)

    card(s, 7.5, 1.5, 5.2, 5.2, bg_color=WHITE, border_color=LIGHT_GRAY)
    tb(s, 7.8, 1.65, 4.6, 0.4, "Reading the Heatmap", size=18, bold=True, color=NAVY)
    bullets(s, 7.8, 2.2, 4.6, 4.2, [
        ("Jaccard = |intersection| / |union|", True, 0),
        ("MIQP variants agree strongly\nwith each other (J = 0.72)", False, 0),
        ("Backward and MIQP are quite\nsimilar (J = 0.72) - they rank\nsite importance similarly", False, 0),
        ("Greedy+Local is most distinct\n(J = 0.35-0.56) due to local\nsearch swaps", False, 0),
        ("Despite different selections, all\nmethods share a consensus core", True, 0),
    ], size=13, color=DARK_GRAY, space=Pt(10))


def slide_consensus(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Consensus Core: High-Confidence Restoration Targets")

    img_centered(s, FIGURES_DIR / "fig4_consensus_sites.png", 1.2, max_w=11.5, max_h=3.6)

    # Three insight cards
    card(s, 0.8, 5.1, 3.7, 2.0, bg_color=RGBColor(0xE8, 0xF5, 0xE9), border_color=ACCENT_GREEN)
    tb(s, 1.0, 5.2, 3.3, 0.3, "All 5 Methods (10 sites)", size=14, bold=True, color=ACCENT_GREEN)
    tb(s, 1.0, 5.55, 3.3, 1.3,
       "10, 15, 31, 32, 37, 40, 41,\n49, 51, 52, 53\n\nHighest-confidence targets\nregardless of methodology",
       size=12, color=DARK_GRAY, line_spacing=Pt(16))

    card(s, 4.8, 5.1, 3.7, 2.0, bg_color=RGBColor(0xFF, 0xF8, 0xE1), border_color=RGBColor(0xFF, 0xC1, 0x07))
    tb(s, 5.0, 5.2, 3.3, 0.3, "4/5 Methods (7 sites)", size=14, bold=True, color=RGBColor(0xF5, 0x7F, 0x17))
    tb(s, 5.0, 5.55, 3.3, 1.3,
       "16, 17, 21, 27, 33, 36, 47, 59\n\nStrong candidates with\nbroad algorithmic support",
       size=12, color=DARK_GRAY, line_spacing=Pt(16))

    card(s, 8.8, 5.1, 3.7, 2.0, bg_color=RGBColor(0xFC, 0xE4, 0xEC), border_color=RGBColor(0xE0, 0x60, 0x60))
    tb(s, 9.0, 5.2, 3.3, 0.3, "1-2 Methods (8 sites)", size=14, bold=True, color=RGBColor(0xC6, 0x28, 0x28))
    tb(s, 9.0, 5.55, 3.3, 1.3,
       "4, 6, 11, 19, 24, 30, 38, etc.\n\nContested sites where\nmethodology choice matters",
       size=12, color=DARK_GRAY, line_spacing=Pt(16))


def slide_network(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Network Analysis: Why Are These Sites Selected?")

    img_centered(s, FIGURES_DIR / "fig5_network_centrality.png", 1.2, max_w=11.5, max_h=4.0)

    card(s, 0.8, 5.5, 11.7, 1.6, bg_color=ICE_BLUE, border_color=BLUE)
    tb(s, 1.1, 5.6, 11.0, 0.3, "Key Insight", size=16, bold=True, color=NAVY)
    tb(s, 1.1, 5.95, 11.0, 1.0,
       "MIQP-selected sites (blue) dominate all network centrality metrics - they are the hubs "
       "that receive and export the most larvae. The optimization naturally discovers network-central sites. "
       "Top hubs (sites 10, 31, 32, 40, 41) are consensus picks across every method tested.",
       size=14, color=DARK_GRAY, line_spacing=Pt(20))


def slide_reef_sizes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "MIQP Extension: Optimal Reef Size Allocation")

    img_centered(s, FIGURES_DIR / "fig6_reef_sizes.png", 1.2, max_w=11.5, max_h=4.2)

    tb(s, 0.8, 5.8, 11.7, 1.2,
       "The MIQP framework naturally extends to variable reef sizing under a total area budget. "
       "Network hubs get maximum allocation (50 units). Peripheral sites get minimal allocation "
       "- included for connectivity value, not local productivity. "
       "Community constraints force inclusion of smaller sites (42, 47, 48) at minimum size.",
       size=13, color=MID_GRAY, line_spacing=Pt(20))


def slide_selection_matrix(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Complete Site Selection Matrix")

    img_centered(s, FIGURES_DIR / "fig3_site_selection_matrix.png", 1.3, max_w=12.0, max_h=4.0)

    tb(s, 0.8, 5.7, 11.7, 1.2,
       "Blue = selected. Count row shows consensus level. Dense blue columns "
       "(10, 15, 31, 32, 40, 41, 49, 51-53) form the indisputable core. "
       "Greedy+Local uniquely selects sites 6, 24, 30, 38, 39, 55, 60.",
       size=13, color=MID_GRAY, align=PP_ALIGN.CENTER, line_spacing=Pt(20))


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    slide_header(s, "Summary of Results")

    # Build table
    rows, cols = 6, 5
    tbl_shape = s.shapes.add_table(rows, cols,
        Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.8))
    tbl = tbl_shape.table

    # Column widths
    col_widths = [2.4, 2.0, 1.8, 2.2, 3.3]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    headers = ["Method", "ODE Score", "Time", "Overlap w/ Best", "Key Trait"]
    data = [
        ["Greedy + Local Search", "1.862  (best)", "32 min", "-", "Best true ODE score"],
        ["Backward Elimination", "1.819  (-2.3%)", "4.7 min", "18 / 25", "Fast, strong quality"],
        ["MIQP (plain)", "1.713  (-8.0%)", "0.32 sec", "14 / 25", "Near-instant, exact surrogate"],
        ["MIQP + Communities", "1.383  (-25.7%)", "< 1 sec", "15 / 25", "Geographic fairness"],
        ["MIQP + Comm + Sizing", "-", "< 1 sec", "13 / 25", "Full policy model"],
    ]

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    row_colors = [ICE_BLUE, WHITE, ICE_BLUE, WHITE, ICE_BLUE]
    for i, (row_data, bg) in enumerate(zip(data, row_colors)):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Calibri"
                p.font.color.rgb = CHARCOAL
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Key findings below table
    card(s, 0.8, 4.5, 11.7, 2.6, bg_color=WHITE, border_color=BLUE)
    tb(s, 1.2, 4.6, 11.0, 0.4, "Key Findings", size=20, bold=True, color=NAVY)

    findings = [
        ("1.", "Greedy+Local achieves the best ODE biomass but requires 32 minutes of compute"),
        ("2.", "MIQP solves in 0.3 seconds - provably optimal under the quadratic surrogate"),
        ("3.", "The 8% gap between MIQP and Greedy quantifies the surrogate approximation error"),
        ("4.", "All methods agree on a consensus core of 10-11 high-value sites"),
        ("5.", "The MIQP framework naturally extends to community fairness and reef sizing constraints"),
    ]

    for i, (num, text) in enumerate(findings):
        y = 5.1 + i * 0.35
        tb(s, 1.3, y, 0.4, 0.3, num, size=13, bold=True, color=BLUE)
        tb(s, 1.7, y, 10.0, 0.3, text, size=13, color=DARK_GRAY)


def slide_takeaways(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_shape(s, 0, 0, Inches(0.15), SLIDE_H, BLUE)

    tb(s, 1.5, 0.5, 10.0, 0.7, "Key Takeaways", size=36, bold=True, color=WHITE, font="Calibri Light")
    add_line(s, Inches(1.5), Inches(1.15), Inches(2.0), BLUE, thickness=Pt(3))

    # Takeaway cards on dark bg
    card_data = [
        ("MIQP surrogate = fast, exact benchmark",
         "Solves in 0.3 seconds what heuristics need 32 minutes for.\n"
         "Enables rapid exploration of constraints and scenarios."),
        ("Heuristics outperform the surrogate on true ODE",
         "Greedy+Local achieves 8% better ODE score than MIQP.\n"
         "The surrogate approximation is the limiting factor, not the solver."),
        ("All methods converge on a consensus core",
         "Sites 10, 15, 31, 32, 37, 40, 41, 49, 51, 52, 53\n"
         "are robust picks regardless of methodology. These are network hubs."),
        ("MIQP framework is extensible to real policy",
         "Community fairness, reef sizing, budget constraints integrate naturally.\n"
         "Quantifies the cost of policy: geographic fairness costs ~26% biomass."),
    ]

    for i, (title, body) in enumerate(card_data):
        y = 1.5 + i * 1.45
        card(s, 1.5, y, 10.3, 1.3, bg_color=DEEP_BLUE, border_color=RGBColor(0x2A, 0x3F, 0x6A))
        tb(s, 1.9, y + 0.1, 9.5, 0.35, title, size=16, bold=True, color=LIGHT_BLUE)
        tb(s, 1.9, y + 0.5, 9.5, 0.7, body, size=13, color=RGBColor(0xB0, 0xBE, 0xC5), line_spacing=Pt(18))


def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    add_shape(s, 0, 0, Inches(0.15), SLIDE_H, BLUE)

    tb(s, 1.5, 1.5, 10.5, 1.0,
       "Future Work", size=32, bold=True, color=WHITE, font="Calibri Light")
    add_line(s, Inches(1.5), Inches(2.2), Inches(2.0), BLUE, thickness=Pt(3))

    bullets(s, 1.8, 2.5, 10.0, 3.0, [
        ("Improve surrogate approximation (per-site A* instead of median)", False, 0),
        ("Robustness analysis across connectivity matrices and parameters", False, 0),
        ("Parallelize ODE evaluations for faster heuristic convergence", False, 0),
        ("Extend to multi-objective formulations (biomass + biodiversity)", False, 0),
    ], size=15, color=RGBColor(0xB0, 0xBE, 0xC5), space=Pt(14))

    # Divider
    add_shape(s, Inches(1.5), Inches(4.8), Inches(10.3), Pt(1), RGBColor(0x2A, 0x3F, 0x6A))

    tb(s, 1.5, 5.2, 10.5, 1.0,
       "Thank You", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")
    tb(s, 1.5, 6.1, 10.5, 0.5,
       "Questions?", size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)              # 1
    slide_problem(prs)            # 2
    slide_model(prs)              # 3
    slide_connectivity(prs)       # 4
    slide_contribution(prs)       # 5
    slide_greedy_curve(prs)       # 6
    slide_scores(prs)             # 7
    slide_timing(prs)             # 8
    slide_gap(prs)                # 9
    slide_overlap(prs)            # 10
    slide_consensus(prs)          # 11
    slide_network(prs)            # 12
    slide_reef_sizes(prs)         # 13
    slide_selection_matrix(prs)   # 14
    slide_summary(prs)            # 15
    slide_takeaways(prs)          # 16
    slide_end(prs)                # 17

    prs.save(str(OUT_PATH))
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
